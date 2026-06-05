#!/usr/bin/env python3
"""Weekly sync PPT generation."""

import io
from pathlib import Path


def extract_template_colors(template_path: Path) -> dict:
    """Read bg / title / body / accent colors from uploaded template."""
    import zipfile
    import xml.etree.ElementTree as ET

    def hx2(h):
        h = h.lstrip('#')
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    result = dict(bg=(64, 49, 82), title=(255, 255, 255), body=(255, 255, 255),
                  text=(255, 255, 255), accent=(148, 182, 210), dim=(144, 144, 168))

    if not template_path or not template_path.exists():
        return result
    try:
        with zipfile.ZipFile(str(template_path)) as z:
            ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'

            theme_files = sorted(f for f in z.namelist() if 'ppt/theme/theme' in f)
            if theme_files:
                tree = ET.fromstring(z.read(theme_files[0]))
                scheme = {}
                for elem in tree.iter():
                    tag = elem.tag.split('}')[-1]
                    if tag in ('lt1', 'lt2', 'dk1', 'dk2', 'accent1', 'accent2'):
                        for child in elem:
                            ctag = child.tag.split('}')[-1]
                            if ctag == 'srgbClr':
                                val = child.get('val', '').lstrip('#')
                            elif ctag == 'sysClr':
                                val = child.get('lastClr', '').lstrip('#')
                            else:
                                val = ''
                            if len(val) == 6:
                                scheme[tag] = hx2(val)
                if 'lt1'    in scheme: result['title']  = scheme['lt1']
                if 'lt2'    in scheme: result['body']   = scheme['lt2']
                if 'dk1'    in scheme: result['text']   = scheme['dk1']
                if 'accent1' in scheme: result['accent'] = scheme['accent1']

            mtree = ET.fromstring(z.read('ppt/slideMasters/slideMaster1.xml'))
            for sf in mtree.iter(f'{{{ns}}}solidFill'):
                for child in sf:
                    if child.tag.split('}')[-1] == 'srgbClr':
                        val = child.get('val', '')
                        if len(val) == 6:
                            result['bg'] = hx2(val)
                            break
                else:
                    continue
                break
    except Exception:
        pass
    return result


def generate_weekly_ppt(config: dict, records: list, template_path: Path = None) -> bytes:
    """Generate weekly sync PPT and return raw bytes."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import PP_PLACEHOLDER as PPH

    R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    using_template = bool(template_path and template_path.exists())

    # ── Setup presentation ───────────────────────────────────────
    if using_template:
        prs = Presentation(str(template_path))
        existing = list(prs.slides)

        cover_layout   = existing[0].slide_layout if existing else prs.slide_masters[0].slide_layouts[0]
        content_layout = existing[2].slide_layout if len(existing) > 2 else (
                         existing[1].slide_layout if len(existing) > 1 else cover_layout)

        sldIdLst = prs.slides._sldIdLst
        for sId in list(sldIdLst):
            rId = sId.get(f'{{{R_NS}}}id')
            prs.part.drop_rel(rId)
            sldIdLst.remove(sId)

        blank_layout = content_layout
    else:
        def _hx(h):
            h = h.lstrip('#')
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank_layout = cover_layout = content_layout = prs.slide_layouts[6]
        _fb = dict(bg=_hx('403152'), title=_hx('FFFFFF'), body=_hx('D0CDD8'),
                   accent=_hx('94B6D2'), dim=_hx('9090A8'))

    # ── Helpers ──────────────────────────────────────────────────
    BODY_TYPES = (PPH.BODY, PPH.OBJECT, PPH.PICTURE)

    def get_ph(slide, *types):
        for ph in slide.placeholders:
            if ph.placeholder_format.type in types:
                return ph
        return None

    def get_body(slide):
        return get_ph(slide, *BODY_TYPES)

    def set_ph(slide, text, *types):
        ph = get_ph(slide, *types)
        if ph:
            ph.text = text
        return ph

    def add_text_fb(slide, text, l, t, w, h, size, color, bold=False, italic=False,
                    align=PP_ALIGN.LEFT, name='Aptos'):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = name
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic

    def add_image_fitted(slide, img_path, x, y, max_w, max_h):
        try:
            pic = slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(max_w))
            lim = Inches(max_h)
            if pic.height > lim:
                ratio = lim / pic.height
                pic.width = int(pic.width * ratio)
                pic.height = lim
        except Exception:
            pass

    # ── Sort records by handler name ─────────────────────────────
    def _key(r):
        h = r.get('handlers', [])
        return h[0]['name'] if h else '\xff'
    records = sorted(records, key=_key)

    presenters = config.get('presenters', '')
    title_text = config.get('title', 'Weekly Report')

    # ── Build project groups ─────────────────────────────────────
    seen_projs, proj_map = [], {}
    for rec in records:
        proj = rec.get('project', '') or ''
        if proj not in proj_map:
            proj_map[proj] = []
            seen_projs.append(proj)
        proj_map[proj].append(rec)

    # ══ TEMPLATE PATH ════════════════════════════════════════════
    if using_template:
        tc = extract_template_colors(template_path)
        txt_color = RGBColor(*tc['text'])

        def add_desc_text(slide, text, y_inch, h_inch=0.7, font_size=None):
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(y_inch), Inches(11.73), Inches(h_inch))
            tf = tb.text_frame
            tf.word_wrap = True
            for i, line in enumerate(text.split('\n')):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = line
                run.font.color.rgb = txt_color
                if font_size:
                    run.font.size = Pt(font_size)

        # Cover
        s1 = prs.slides.add_slide(cover_layout)
        set_ph(s1, title_text, PPH.CENTER_TITLE, PPH.TITLE)
        if presenters:
            set_ph(s1, f'Presented by {presenters}', PPH.SUBTITLE, PPH.BODY)

        # TOC
        s2 = prs.slides.add_slide(content_layout)
        set_ph(s2, 'Content', PPH.TITLE, PPH.CENTER_TITLE)
        body_ph = get_body(s2)
        if body_ph:
            tf = body_ph.text_frame
            tf.clear()
            first = True
            for proj in seen_projs:
                if proj:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = proj
                    p.level = 0
                    if p.runs: p.runs[0].font.bold = True
                for rec in proj_map[proj]:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    p.text = rec.get('taskText', '')
                    p.level = 1 if proj else 0

        # Content slides
        for rec in records:
            task_text = rec.get('taskText', '')
            notes = rec.get('notes', '').strip()
            valid_imgs = [img for img in rec.get('images', [])
                          if Path(img.get('path', '')).exists()]

            if not valid_imgs:
                slide = prs.slides.add_slide(content_layout)
                set_ph(slide, task_text, PPH.TITLE, PPH.CENTER_TITLE)
                if notes:
                    sn = prs.slides.add_slide(content_layout)
                    set_ph(sn, task_text, PPH.TITLE, PPH.CENTER_TITLE)
                    add_desc_text(sn, notes, 1.55, 5.0, font_size=28)
            else:
                if notes:
                    sn = prs.slides.add_slide(content_layout)
                    set_ph(sn, task_text, PPH.TITLE, PPH.CENTER_TITLE)
                    add_desc_text(sn, notes, 1.55, 5.0, font_size=28)
                for i, img in enumerate(valid_imgs):
                    slide = prs.slides.add_slide(content_layout)
                    set_ph(slide, task_text, PPH.TITLE, PPH.CENTER_TITLE)
                    body = get_body(slide)
                    if body: body.text = ''

                    caption = img.get('caption', '').strip()
                    if caption:
                        desc_h = min(0.45 * len(caption.split('\n')), 1.2)
                        add_desc_text(slide, caption, 1.47, desc_h)
                        img_top = 1.47 + desc_h + 0.1
                    else:
                        img_top = 1.55
                    add_image_fitted(slide, img['path'], 0.8, img_top, 11.73, 6.5 - img_top)

    # ══ FALLBACK (no template) ════════════════════════════════════
    else:
        BG, WHITE, LIGHT, ACCENT, DIM = (_fb['bg'], _fb['title'], _fb['body'],
                                          _fb['accent'], _fb['dim'])

        def set_bg(slide):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = BG

        def make_content(task_text):
            slide = prs.slides.add_slide(blank_layout)
            set_bg(slide)
            add_text_fb(slide, task_text, 0.8, 0.40, 11.73, 0.90, 28, WHITE,
                        bold=True, name='Aptos Display')
            return slide

        s1 = prs.slides.add_slide(blank_layout)
        set_bg(s1)
        add_text_fb(s1, title_text, 1.5, 2.5, 10.33, 1.5, 40, WHITE,
                    bold=True, align=PP_ALIGN.CENTER, name='Aptos Display')
        if presenters:
            add_text_fb(s1, f'Presented by {presenters}', 1.5, 4.2, 10.33, 0.8, 18, LIGHT,
                        align=PP_ALIGN.CENTER)

        s2 = prs.slides.add_slide(blank_layout)
        set_bg(s2)
        add_text_fb(s2, 'Content', 0.8, 0.40, 11.73, 0.90, 28, WHITE, bold=True, name='Aptos Display')
        y = 1.55
        for proj in seen_projs:
            if proj:
                if y + 0.44 > 6.5: break
                add_text_fb(s2, proj, 0.8, y, 11.73, 0.44, 16, ACCENT, bold=True)
                y += 0.44
            for rec in proj_map[proj]:
                if y + 0.40 > 6.5: break
                indent = 1.3 if proj else 0.8
                add_text_fb(s2, rec.get('taskText', ''), indent, y, 11.73 - (indent - 0.8), 0.40, 14, LIGHT)
                y += 0.40
        add_text_fb(s2, '1', 12.3, 6.9, 1, 0.4, 11, DIM, align=PP_ALIGN.RIGHT)

        slide_num = 2
        for rec in records:
            task_text = rec.get('taskText', '')
            notes = rec.get('notes', '').strip()
            valid_imgs = [img for img in rec.get('images', [])
                          if Path(img.get('path', '')).exists()]
            if not valid_imgs:
                slide = make_content(task_text)
                add_text_fb(slide, str(slide_num), 12.3, 6.9, 1, 0.4, 11, DIM, align=PP_ALIGN.RIGHT)
                slide_num += 1
                if notes:
                    slide = make_content(task_text)
                    add_text_fb(slide, notes, 0.8, 1.55, 11.73, 5.0, 28, LIGHT)
                    add_text_fb(slide, str(slide_num), 12.3, 6.9, 1, 0.4, 11, DIM, align=PP_ALIGN.RIGHT)
                    slide_num += 1
            else:
                if notes:
                    slide = make_content(task_text)
                    add_text_fb(slide, notes, 0.8, 1.55, 11.73, 5.0, 28, LIGHT)
                    add_text_fb(slide, str(slide_num), 12.3, 6.9, 1, 0.4, 11, DIM, align=PP_ALIGN.RIGHT)
                    slide_num += 1
                for i, img in enumerate(valid_imgs):
                    slide = make_content(task_text)
                    caption = img.get('caption', '').strip()
                    if caption:
                        add_text_fb(slide, caption, 0.8, 1.47, 11.73, 0.70, 20, LIGHT)
                        img_top = 2.25
                    else:
                        img_top = 1.55
                    add_image_fitted(slide, img['path'], 0.8, img_top, 11.73, 6.5 - img_top)
                    add_text_fb(slide, str(slide_num), 12.3, 6.9, 1, 0.4, 11, DIM, align=PP_ALIGN.RIGHT)
                    slide_num += 1

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
